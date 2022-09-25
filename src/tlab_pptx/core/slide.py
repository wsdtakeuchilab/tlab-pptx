import dataclasses
import io

import plotly.graph_objects as go
import pptx
import pptx.shapes.autoshape
import pptx.shapes.placeholder
import pptx.shapes.shapetree
import pptx.slide
import pptx.util


@dataclasses.dataclass()
class Slide:
    """
    A wrapper class for pptx.slide.Slide.
    """

    _slide: pptx.slide.Slide
    """The internal slide."""

    def update_title(
        self,
        text: str | None = None,
        font_name: str | None = None,
        font_size: int | None = None,
        font_bold: bool | None = None,
        font_italic: bool | None = None,
    ) -> "Slide":
        """
        Updates the title of the slide.

        Parameters
        ----------
        title_text : str
            The title text.
        font_name : str
            The font name of the title text.
        font_size : int
            The size of the title text in point.
        font_bold : bool
            If true, the title text is bold style.
        font_italic : bool
            If true, the title text is italic style.

        Returns
        -------
        tlab_pptx.core.slide.Slide
            Itself.

        Raises
        ------
        ValueError
            If the slide does not have a title shape.
        """
        title = self._slide.shapes.title
        if title is None:
            raise ValueError("No title placeholder in this slide")
        assert isinstance(title, pptx.shapes.placeholder.SlidePlaceholder)
        if text is not None:
            title.text = text
        for paragraph in title.text_frame.paragraphs:
            if font_name is not None:
                paragraph.font.name = font_name
            if font_size is not None:
                paragraph.font.size = pptx.util.Pt(font_size)
            if font_bold is not None:
                paragraph.font.bold = font_bold
            if font_italic is not None:
                paragraph.font.italic = font_italic
        return self

    def add_text(
        self,
        text: str,
        left: float,
        top: float,
        width: float = 6.0,
        height: float = 4.0,
        font_name: str = "Arial",
        font_size: int = 18,
        font_bold: bool = False,
        font_italic: bool = False,
    ) -> "Slide":
        """
        Adds a text to the slide.

        Parameters
        ----------
        text : str
            A text to be added.
        left : float
            The left position of the text in centimeter.
        top : float
            The top position of the text in centimeter.
        width : float
            The width of the text in centimeter.
        height : float
            The height of the text in centimeter.
        font_name : str
            The font name of the text.
        font_size : int
            The size of the text in point.
        font_bold : bool
            If true, the text is bold style.
        font_italic : bool
            If true, the text is italic style.

        Returns
        -------
        tlab_pptx.core.slide.Slide
            Itself.
        """
        shapes = self._slide.shapes
        assert isinstance(shapes, pptx.shapes.shapetree.SlideShapes)
        textbox = shapes.add_textbox(
            left=pptx.util.Cm(left),
            top=pptx.util.Cm(top),
            width=pptx.util.Cm(width),
            height=pptx.util.Cm(height),
        )
        assert isinstance(textbox, pptx.shapes.autoshape.Shape)
        textbox.text_frame.text = text
        for paragraph in textbox.text_frame.paragraphs:
            paragraph.font.name = font_name
            paragraph.font.size = pptx.util.Pt(font_size)
            paragraph.font.bold = font_bold
            paragraph.font.italic = font_italic
        return self

    def add_figure(
        self,
        fig: go.Figure,
        left: float,
        top: float,
        width: float = 11.5,
        height: float = 11.5,
    ) -> "Slide":
        """
        Adds a figure to the slide.

        Parameters
        ----------
        fig : plotly.graph_objects.Figure
            A figure to be added.
        left : float
            The left position of the figure in centimeter.
        top : float
            The top position of the figure in centimeter.
        width : float
            The width of the figure in centimeter.
        height : float
            The height of the figure in centimeter.

        Returns
        -------
        tlab_pptx.core.slide.Slide
            Itself.
        """
        shapes = self._slide.shapes
        assert isinstance(shapes, pptx.shapes.shapetree.SlideShapes)
        with io.BytesIO(fig.to_image("png", scale=10)) as f:
            shapes.add_picture(
                f,
                left=pptx.util.Cm(left),
                top=pptx.util.Cm(top),
                width=pptx.util.Cm(width),
                height=pptx.util.Cm(height),
            )
        return self
