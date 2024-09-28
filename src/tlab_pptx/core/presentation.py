import dataclasses
import typing as t

import pptx
import pptx.presentation
import pptx.slide

from tlab_pptx import pptx as tpptx
from tlab_pptx.core import slide


@dataclasses.dataclass()
class Presentation:
    """
    A wrapper class for pptx.presentation.Presentation.
    """

    _prs: pptx.presentation.Presentation
    """The internal presentation."""

    @property
    def slide_layouts(self) -> tuple[pptx.slide.SlideLayout, ...]:
        """
        A tuple of slide layouts belonging to the first slide master.
        """
        return tuple(self._prs.slide_layouts)

    @property
    def slides(self) -> tuple[slide.Slide, ...]:
        """
        A tuple of slides in the presentation.
        """
        return tuple(map(slide.Slide, self._prs.slides))

    def add_slide(self, layout_idx: int | None = None) -> slide.Slide:
        """
        Creates a new slide and add it to the presentation.

        Parameters
        ----------
        layout_idx : int
            An index of Presentation.slide_layouts to be inherited by the slide.

        Returns
        -------
        tlab_pptx.core.slide.Slide
            A new slide.
        """
        if layout_idx is None:
            layout_idx = 0
        slides = self._prs.slides
        sld = slides.add_slide(self.slide_layouts[layout_idx])
        return slide.Slide(sld)

    def save(self, file: str | t.IO[bytes]) -> None:
        """
        Saves as a `pptx` file.

        Parameters
        ----------
        file : str | IO[bytes]
            A filepath string or buffer object to which the presentation is saved.
        """
        self._prs.save(file)


def new_presentation(
    file: str | t.IO[bytes] | None = None,
) -> Presentation:
    """
    Creates a new presentation.

    Parameters
    ----------
    file : tlab_pptx.str | IO[bytes] | None
        A filepath string or buffer object of a `.pptx` file as a base.
        if None (default), `tlab_pptx.pptx.DEFAULT_PPTX` is used.

    Returns
    -------
    tlab_pptx.core.presentation.Presentation
        A new presentation.
    """
    if file is None:
        file = tpptx.DEFAULT_PPTX
    prs = pptx.Presentation(file)
    return Presentation(prs)
