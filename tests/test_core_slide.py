import copy
from collections import abc
from unittest import mock

import plotly.graph_objects as go
import pptx
import pptx.shapes.autoshape
import pptx.shapes.placeholder
import pptx.shapes.shapetree
import pptx.slide
import pptx.text.text
import pptx.util
import pytest
from tlab_pptx.core import slide as tslide


def assert_add_picture_called_with(
    slide: tslide.Slide,
    left: float | None = None,
    top: float | None = None,
    width: float | None = None,
    height: float | None = None,
) -> None:
    slide._slide.shapes.add_picture.assert_called_once_with(
        mock.ANY,
        left=pptx.util.Cm(left) if left is not None else mock.ANY,
        top=pptx.util.Cm(top) if top is not None else mock.ANY,
        width=pptx.util.Cm(width) if width is not None else mock.ANY,
        height=pptx.util.Cm(height) if height is not None else mock.ANY,
    )


def assert_add_textbox_called_with(
    slide: tslide.Slide,
    left: float | None = None,
    top: float | None = None,
    width: float | None = None,
    height: float | None = None,
) -> None:
    slide._slide.shapes.add_textbox.assert_called_once_with(
        left=pptx.util.Cm(left) if left is not None else mock.ANY,
        top=pptx.util.Cm(top) if top is not None else mock.ANY,
        width=pptx.util.Cm(width) if width is not None else mock.ANY,
        height=pptx.util.Cm(height) if height is not None else mock.ANY,
    )


def describe_slide() -> None:
    @pytest.fixture()
    def slide() -> tslide.Slide:
        slide_mock = mock.Mock(spec_set=pptx.slide.Slide)
        slide_mock.shapes = mock.Mock(spec_set=pptx.shapes.shapetree.SlideShapes)
        slide_mock.shapes.title = mock.Mock(
            spec_set=pptx.shapes.placeholder.SlidePlaceholder
        )
        slide_mock.shapes.title.text_frame.paragraphs = [
            mock.Mock(spec_set=pptx.text.text._Paragraph)
        ]
        slide_mock.shapes.add_textbox.return_value = mock.Mock(
            spec_set=pptx.shapes.autoshape.Shape
        )
        slide_mock.shapes.add_textbox.return_value.text_frame.paragraphs = [
            mock.Mock(spec_set=pptx.text.text._Paragraph)
        ]
        return tslide.Slide(slide_mock)

    @pytest.mark.parametrize("text", [None, "hello", "hello\ngoodbye"])
    def test_update_title_text(
        slide: tslide.Slide,
        text: str | None,
    ) -> None:
        title = slide._slide.shapes.title
        title_copy = copy.copy(title)
        slide.update_title(text=text)
        if text is None:
            assert title.text == title_copy.text
        else:
            assert title.text == text

    @pytest.mark.parametrize("font_name", [None, "Arial", "San Serif"])
    def test_update_title_font_name(slide: tslide.Slide, font_name: str | None) -> None:
        title = slide._slide.shapes.title
        title_copy = copy.copy(title)
        slide.update_title(font_name=font_name)
        for i, paragraph in enumerate(title.text_frame.paragraphs):
            font = paragraph.font
            if font_name is None:
                font_copy = title_copy.text_frame.paragraphs[i].font
                assert font.name == font_copy.name
            else:
                assert font.name == font_name

    @pytest.mark.parametrize("font_size", [None, 18, 28])
    def test_update_title_font_size(slide: tslide.Slide, font_size: int | None) -> None:
        title = slide._slide.shapes.title
        title_copy = copy.copy(title)
        slide.update_title(font_size=font_size)
        for i, paragraph in enumerate(title.text_frame.paragraphs):
            font = paragraph.font
            if font_size is None:
                font_copy = title_copy.text_frame.paragraphs[i].font
                assert font.size == font_copy.size
            else:
                assert font.size == pptx.util.Pt(font_size)

    @pytest.mark.parametrize("font_bold", [None, True, False])
    def test_update_title_font_bold(
        slide: tslide.Slide, font_bold: bool | None
    ) -> None:
        title = slide._slide.shapes.title
        title_copy = copy.copy(title)
        slide.update_title(font_bold=font_bold)
        for i, paragraph in enumerate(title.text_frame.paragraphs):
            font = paragraph.font
            if font_bold is None:
                font_copy = title_copy.text_frame.paragraphs[i].font
                assert font.bold == font_copy.bold
            else:
                assert font.bold == font_bold

    @pytest.mark.parametrize("font_italic", [None, True, False])
    def test_update_title_font_italic(
        slide: tslide.Slide, font_italic: bool | None
    ) -> None:
        title = slide._slide.shapes.title
        title_copy = copy.copy(title)
        slide.update_title(font_italic=font_italic)
        for i, paragraph in enumerate(title.text_frame.paragraphs):
            font = paragraph.font
            if font_italic is None:
                font_copy = title_copy.text_frame.paragraphs[i].font
                assert font.italic == font_copy.italic
            else:
                assert font.italic == font_italic

    def test_update_title_no_title(slide: tslide.Slide) -> None:
        slide._slide.shapes.title = None
        with pytest.raises(ValueError):
            slide.update_title()

    @pytest.fixture()
    def fig() -> abc.Generator[go.Figure, None, None]:
        with mock.patch("plotly.graph_objects.Figure.to_image", return_value=b""):
            yield go.Figure()

    @pytest.mark.parametrize("left", [0.0, 1.0])
    @pytest.mark.parametrize("top", [0.0, 1.0])
    def test_add_figure(
        slide: tslide.Slide,
        fig: go.Figure,
        left: float,
        top: float,
    ) -> None:
        assert slide.add_figure(fig, left, top) == slide
        fig.to_image.assert_called_once_with("png", scale=10)
        assert_add_picture_called_with(slide, left, top)

    @pytest.mark.parametrize("width", [0.0, 1.0])
    def test_add_figure_size(
        slide: tslide.Slide,
        fig: go.Figure,
        width: float,
    ) -> None:
        left, top = 0.0, 0.0
        assert slide.add_figure(fig, left, top, width=width) == slide
        assert_add_picture_called_with(slide, left, top, width)

    @pytest.mark.parametrize("height", [0.0, 1.0])
    def test_add_figure_height(
        slide: tslide.Slide,
        fig: go.Figure,
        height: float,
    ) -> None:
        left, top = 0.0, 0.0
        assert slide.add_figure(fig, left, top, height=height) == slide
        assert_add_picture_called_with(slide, left, top, height=height)

    @pytest.fixture()
    def text() -> str:
        return "hello"

    @pytest.mark.parametrize("text", ["hello", "hello\ngoodbye"])
    @pytest.mark.parametrize("left", [0.0, 1.0])
    @pytest.mark.parametrize("top", [0.0, 1.0])
    def test_add_text(
        slide: tslide.Slide,
        text: str,
        left: float,
        top: float,
    ) -> None:
        assert slide.add_text(text, left, top) == slide
        assert_add_textbox_called_with(slide, left, top)

    @pytest.mark.parametrize("width", [0.0, 1.0])
    def test_add_text_width(
        slide: tslide.Slide,
        text: str,
        width: float,
    ) -> None:
        left, top = 0.0, 0.0
        assert slide.add_text(text, left, top, width=width) == slide
        assert_add_textbox_called_with(slide, left, top, width=width)

    @pytest.mark.parametrize("height", [0.0, 1.0])
    def test_add_text_height(
        slide: tslide.Slide,
        text: str,
        height: float,
    ) -> None:
        left, top = 0.0, 0.0
        assert slide.add_text(text, left, top, height=height) == slide
        assert_add_textbox_called_with(slide, left, top, height=height)

    @pytest.mark.parametrize("font_name", ["Arial", "San Serif"])
    def test_add_text_font_name(
        slide: tslide.Slide,
        text: str,
        font_name: str,
    ) -> None:
        slide.add_text(
            text,
            0.0,
            0.0,
            font_name=font_name,
        )
        textbox = slide._slide.shapes.add_textbox.return_value
        for paragraph in textbox.text_frame.paragraphs:
            assert paragraph.font.name == font_name

    @pytest.mark.parametrize("font_size", [18, 28])
    def test_add_text_font_size(
        slide: tslide.Slide,
        text: str,
        font_size: int,
    ) -> None:
        slide.add_text(
            text,
            0.0,
            0.0,
            font_size=font_size,
        )
        textbox = slide._slide.shapes.add_textbox.return_value
        for paragraph in textbox.text_frame.paragraphs:
            assert paragraph.font.size == pptx.util.Pt(font_size)

    @pytest.mark.parametrize("font_bold", [True, False])
    def test_add_text_font_bold(
        slide: tslide.Slide,
        text: str,
        font_bold: bool,
    ) -> None:
        slide.add_text(
            text,
            0.0,
            0.0,
            font_bold=font_bold,
        )
        textbox = slide._slide.shapes.add_textbox.return_value
        for paragraph in textbox.text_frame.paragraphs:
            assert paragraph.font.bold == font_bold

    @pytest.mark.parametrize("font_italic", [True, False])
    def test_add_text_font_italic(
        slide: tslide.Slide,
        text: str,
        font_italic: bool,
    ) -> None:
        slide.add_text(
            text,
            0.0,
            0.0,
            font_italic=font_italic,
        )
        textbox = slide._slide.shapes.add_textbox.return_value
        for paragraph in textbox.text_frame.paragraphs:
            assert paragraph.font.italic == font_italic
