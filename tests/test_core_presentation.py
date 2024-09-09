import functools
from collections import abc
from unittest import mock

import pptx
import pptx.presentation
import pptx.slide
import pytest

from tlab_pptx import pptx as tpptx
from tlab_pptx import typing
from tlab_pptx.core import presentation, slide


def test_new_presentation() -> None:
    with mock.patch("pptx.Presentation", return_value=pptx.Presentation()) as m:
        prs = presentation.new_presentation()
    m.assert_called_once_with(tpptx.DEFAULT_PPTX)
    assert prs._prs == m.return_value


@pytest.mark.parametrize(
    ["filename", "open_mode"],
    [(tpptx.DEFAULT_PPTX, "rb")],
)
def test_new_presentation_filepath_or_buffer(
    filepath_or_buffer: typing.FilePathOrBuffer | None,
) -> None:
    with mock.patch("pptx.Presentation", return_value=pptx.Presentation()) as m:
        prs = presentation.new_presentation(filepath_or_buffer)
    m.assert_called_once_with(filepath_or_buffer)
    assert prs._prs == m.return_value


@functools.lru_cache(maxsize=32)
def _new_presentation(
    filepath: typing.FilePath | None = None,
) -> pptx.presentation.Presentation:
    return pptx.Presentation(filepath)


def describe_presentation() -> None:
    @pytest.fixture()
    def prs() -> presentation.Presentation:
        return presentation.Presentation(_new_presentation())

    def test_slides(prs: presentation.Presentation) -> None:
        assert prs.slides == tuple(map(slide.Slide, prs._prs.slides))

    @pytest.fixture()
    def add_slide_mock() -> abc.Generator[mock.Mock, None, None]:
        sld = mock.Mock(spec_set=pptx.slide.Slide)
        with mock.patch("pptx.slide.Slides.add_slide", return_value=sld) as m:
            yield m

    def test_add_slide(
        add_slide_mock: mock.Mock, prs: presentation.Presentation
    ) -> None:
        assert prs.add_slide() == slide.Slide(add_slide_mock.return_value)
        add_slide_mock.assert_called_once_with(prs._prs.slide_layouts[0])

    @pytest.mark.parametrize("layout_idx", [0, 1])
    def test_add_slide_layout_idx(
        add_slide_mock: mock.Mock, prs: presentation.Presentation, layout_idx: int
    ) -> None:
        assert prs.add_slide(layout_idx) == slide.Slide(add_slide_mock.return_value)
        add_slide_mock.assert_called_once_with(prs._prs.slide_layouts[layout_idx])

    @pytest.mark.parametrize(
        ["filename", "open_mode"],
        [("test_presentation_save.pptx", "wb")],
    )
    @mock.patch("pptx.presentation.Presentation.save")
    def test_save(
        save_mock: mock.Mock,
        prs: presentation.Presentation,
        filepath_or_buffer: typing.FilePathOrBuffer,
    ) -> None:
        prs.save(filepath_or_buffer)
        save_mock.assert_called_once_with(filepath_or_buffer)
